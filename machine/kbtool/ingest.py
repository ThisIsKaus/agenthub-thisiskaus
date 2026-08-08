#!/usr/bin/env python3
"""
AgentHub knowledge base ingest.

Structure-aware chunking, sensitivity classification, content-hash deduplication,
and hard exclusion of employer systems, credential stores and prior-client artefacts.
Every chunk carries a sensitivity class so retrieval can decide what a cloud lane may see.
Nothing classed S1c, S2 or S3 leaves this machine.

  ingest.py                  incremental
  ingest.py --rebuild        drop and rebuild (required after any classifier change)
  ingest.py --dry            classify and report, write nothing
  ingest.py --only S1p       one class at a time
  ingest.py --list S3        print every path in a class and stop
  ingest.py --blocked        print what was excluded and why
"""

import argparse, hashlib, html as _html, json, os, re, sys, time
from pathlib import Path

import requests
import lancedb

HOME = Path.home()
H = HOME / "AgentHub"
DB_PATH = H / "kb"
TABLE = "kb_main"
EMBED_URL = "http://127.0.0.1:4000/v1/embeddings"
EMBED_MODEL = "local-embed"
CHUNK, OVERLAP, BATCH = 1600, 200, 32
EXTS = {".md", ".txt", ".pdf", ".docx", ".xlsx", ".pptx", ".csv", ".html", ".htm"}

SOURCES = [H / "canon", H / "inbox", H / "docs", H / "drafts", H / "digests",
           H / "skills-lib" / "skills", H / "state" / "proposals",
           H / "state" / "builds", H / "evals", H / "contracts", H / "logs"]
# Canvas should be able to answer a question about the system itself: which skills cover
# negotiation, what the diagnostician proposed and why, what a build actually changed, what
# the last self-test failed on. Those artefacts existed and were unreadable — the system knew
# things it could not tell you.
# docs/ already carries memory notes and exported sessions; digests join them so a
# question, an answer and the day it arrived are all findable from one query.
EXTRA = H / "kbtool" / "sources.json"

# --------------------------------------------------------------------- blocks
# Employer SYSTEMS and TOOL artefacts. Not the company name in general: Kos's own
# career history lives in folders named after employers, is his, and is classified
# CAREER -> S3 (local only) rather than blocked.
EMPLOYER = re.compile(
    r"(microsoft\s+copilot\s+chat|microsoft\s+file\s+share|microsoft\s+edge\s+drop|"
    r"onedrive\s*-\s*microsoft|@microsoft\.com|corp\.microsoft|\bmsx\b|mspartner|"
    r"microsoft\s+internal|confidential[\s_-]internal)", re.I)

# Credential stores. Never read, never embedded, under any classification.
SECRETS = re.compile(
    r"(1password|lastpass|bitwarden|keepass|emergency[\s_-]?kit|recovery[\s_-]?(key|code|kit)|"
    r"secret[\s_-]?key|master[\s_-]?password|seed[\s_-]?phrase|private[\s_-]?key|"
    r"backup[\s_-]?code|credentials?\b|\bpasswords?\b|passwd|\bpwd\b|\.env\b|"
    r"\.pem\b|\.p12\b|\.key\b|api[\s_-]?key|auth[\s_-]?token)", re.I)

# Delivery artefacts from prior engagements: third-party confidential, no forward
# value, and hundreds of near-identical machine-generated files that poison retrieval.
# Superseded versions sit alongside current ones and win retrieval on lexical similarity
# while carrying stale facts. The current version is already indexed; these are noise.
SUPERSEDED = re.compile(
    r"(superseded|deprecated|obsolete|\(old\)|[\s_-]old[\s_-]version|"
    r"[\s_-]v\d+[\s_-]?\(old|do[\s_-]not[\s_-]use|backup[\s_-]copy|copy[\s_-]of[\s_-])", re.I)

CLIENT_ARTEFACT = re.compile(
    r"(snapshot!|csu[\s_-]uat|adha[\s_-]dev|customizations/viewdetail|"
    r"-privileges-|-viewdetail-|msdyn_|-securityrole-)", re.I)

# ------------------------------------------------------------- classification
PRODUCT = re.compile(
    r"(my[\s_-]?financial[\s_-]?compass|metascan|veritrail|agentcraft|orginsight|"
    r"gencai|openclaw|polyclaw|agenthub|\bapvt\b)", re.I)

S1C = re.compile(
    r"(/clients?/|customer[\s_-]engagement|client[\s_-]engagement|"
    r"vinnies|meesho|\bAGN-[A-Z]{3}-)", re.I)

MONEY = re.compile(
    r"(bank\s*(statement|account|transfer|detail)|commbank|westpac|\banz\b|\bnab\b|hsbc|"
    r"icici|hdfc|citibank|bankwest|\bsbi\b|\bstmt\b|"
    r"statement|payslip|pay[\s_-]?slip|salary|super(annuation)?|smsf|tax\b|itr\b|"
    r"mortgage|home\s*loan|zerodha|stake\b|fidelity|espp|\brsu\b|net[\s_-]?worth|portfolio|"
    r"dividend|invoice|investor|agreement|contract|premium|policy[\s_-]?no|"
    r"deed|conveyanc|settlement|valuation|rates[\s_-]notice|expense[\s_-]sheet)", re.I)

IDENTITY = re.compile(
    r"(passport|visa\b|medicare|driver.{0,3}licen|aadhaa?r|pan[\s_-]?card|\btfn\b|"
    r"birth[\s_-]?cert|marriage[\s_-]?cert|citizen(ship)?|\boci\b|credit[\s_-]?report|"
    r"police[\s_-]?check|medical|pathology|prescription|immunis|vaccin|identity|"
    r"\bwill\b|power[\s_-]of[\s_-]attorney|probate)", re.I)

CAREER = re.compile(
    r"(/work/|resum|\bcv\b|curriculum[\s_-]vitae|performance[\s_-]review|"
    r"promotion|compensation|offer[\s_-]letter|reference[\s_-]letter|"
    r"employment|appraisal|career)", re.I)

S2 = re.compile(r"(envelope[\s_-]?(collective|tales)|theenvelopecollective|neelam)", re.I)

S1P = re.compile(r"(agenticality|nxi\b|agentic[\s_-]personality)", re.I)

SKIP_PARTS = {"node_modules", ".venv", "venv", "site-packages", "__pycache__",
              "Caches", "Application Support", ".git", ".Trash", "obj", "bin",
              "TestResults", "target", "Pods", ".next", ".cache"}


def block_reason(path: Path):
    """Why a file is never read. None means allowed."""
    s = str(path)
    if EMPLOYER.search(s):
        return "employer"
    if SECRETS.search(s):
        return "credentials"
    if CLIENT_ARTEFACT.search(s):
        return "prior-client artefact"
    if SUPERSEDED.search(s):
        return "superseded version"
    return None


def excluded(path: Path) -> bool:
    return block_reason(path) is not None


def classify(path: Path, default="S0") -> str:
    s = str(path)
    # Derived material inherits the sensitivity of what it derives from. A memory note
    # distils a day that may have included S3 questions; a digest carries email subjects
    # across every class. Neither may reach a cloud lane on the strength of its own path.
    if "/docs/memory/" in s or "/digests/" in s or "/sessions/" in s or "/logs/" in s:
        return "S3"
    if "/skills-lib/" in s or "/state/proposals/" in s or "/state/builds/" in s \
            or "/contracts/" in s or "/evals/" in s:
        return "S1p"
    if S1C.search(s):
        return "S1c"
    if PRODUCT.search(s):
        return "S1p"
    if MONEY.search(s) or IDENTITY.search(s) or CAREER.search(s):
        return "S3"
    if S2.search(s):
        return "S2"
    if S1P.search(s):
        return "S1p"
    return default


# -------------------------------------------------------------------- reading
def read_text(p: Path) -> str:
    suf = p.suffix.lower()
    if suf == ".pdf":
        import fitz
        with fitz.open(p) as doc:
            return "\n".join(page.get_text() for page in doc)
    if suf == ".docx":
        import docx
        return "\n".join(x.text for x in docx.Document(str(p)).paragraphs if x.text.strip())
    if suf == ".xlsx":
        import openpyxl
        wb = openpyxl.load_workbook(str(p), read_only=True, data_only=True)
        out = []
        for ws in wb.worksheets:
            out.append(f"## {ws.title}")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
                if cells:
                    out.append(" | ".join(cells))
        wb.close()
        return "\n".join(out)
    if suf == ".pptx":
        from pptx import Presentation
        out = []
        for i, slide in enumerate(Presentation(str(p)).slides, 1):
            out.append(f"## Slide {i}")
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    out.append(shape.text_frame.text.strip())
        return "\n".join(out)
    if suf in (".html", ".htm"):
        raw = p.read_text(errors="ignore")
        raw = re.sub(r"(?is)<(script|style)[^>]*>.*?</\1>", " ", raw)
        return _html.unescape(re.sub(r"\s{2,}", " ", re.sub(r"(?s)<[^>]+>", " ", raw)))
    return p.read_text(errors="ignore")


def context_label(path: Path) -> str:
    """DISABLED — measured 8 Aug and reverted. Recall fell 87% to 81%, MRR 0.751 to 0.701.

    Two faults. The date regex matched inside digit runs: PaymentAdvice_22022019.pdf (22 Feb
    2019) was labelled 2022-01, because "2022" appears inside "22022019". Every dated document
    carried a fabricated date into its embedding, and false metadata is worse than none.
    Second, the longer prefix changed dedup: chunks fell 50,017 to 38,766 and duplicate files
    rose 417 to 1,195, so real content was dropped.

    The folder trail was correct and useful — "OpenClaw / 03-Knowledge-Base" does distinguish
    one README from seventy. Re-enable that half only, with the date read from file mtime or
    a parsed path segment rather than a regex over the whole string.

    A situating line for every chunk: where the file lives, and when it is about.

    The chunker prepended the filename alone, which discriminates nothing when 286 payslips
    share the name KaustubhBajpai_PaySlip.pdf and 70 files are called README.md. Seven of
    twelve golden-set misses were documents identical but for a date the embedding never saw,
    and three wanted one README among seventy.

    The folder trail and any date in the path go into both the vector and the BM25 index.
    Deterministic, no model call — the cheap half of contextual retrieval, aimed at the
    failure the golden set actually reports.
    """
    home = str(Path.home())
    rel = str(path).replace(home + "/", "")
    for noise in ("Library/CloudStorage/", "OneDrive-Personal/",
                  "OneDrive-AgenticPersonalityPty.Ltd/", "AgentHub/"):
        rel = rel.replace(noise, "")
    parts = [x for x in Path(rel).parts[:-1] if x not in (".", "")][-3:]
    m = re.search(r"(20\d{2})[-_ ]?(0[1-9]|1[0-2])?[-_ ]?(0[1-9]|[12]\d|3[01])?", rel)
    when = ""
    if m:
        when = m.group(1) + ("-" + m.group(2) if m.group(2) else "") + \
               ("-" + m.group(3) if m.group(3) else "")
    bits = [b for b in (when, " / ".join(parts), path.name) if b]
    return " · ".join(bits)


def chunks(text, name=""):
    """Split on markdown headings, then paragraphs. Each chunk carries its heading trail."""
    sections, trail, buf = [], [], []

    def flush():
        if buf and any(l.strip() for l in buf):
            sections.append((" / ".join(trail), "\n".join(buf).strip()))

    for line in text.splitlines():
        m = re.match(r"^(#{1,4})\s+(.*)", line)
        if m:
            flush()
            buf.clear()
            depth = len(m.group(1))
            trail = trail[:depth - 1] + [m.group(2).strip()]
        else:
            buf.append(line)
    flush()
    if not sections:
        sections = [("", text)]

    for head, body in sections:
        prefix = (f"{name} — {head}\n" if head else (f"{name}\n" if name else ""))
        body = re.sub(r"\n{3,}", "\n\n", body).strip()
        if not body:
            continue
        if len(body) <= CHUNK:
            yield prefix + body
            continue
        cur = ""
        for para in body.split("\n\n"):
            if cur and len(cur) + len(para) + 2 > CHUNK:
                yield prefix + cur.strip()
                cur = para
            else:
                cur = (cur + "\n\n" + para) if cur else para
            while len(cur) > CHUNK:
                yield prefix + cur[:CHUNK]
                cur = cur[CHUNK - OVERLAP:]
        if cur.strip():
            yield prefix + cur.strip()


def embed(texts):
    payload = {"model": EMBED_MODEL, "input": ["search_document: " + t for t in texts]}
    r = requests.post(EMBED_URL, json=payload, timeout=180)
    r.raise_for_status()
    data = sorted(r.json()["data"], key=lambda d: d["index"])
    return [d["embedding"] for d in data]


def gather():
    roots = [(p, "S0") for p in SOURCES]
    if EXTRA.exists():
        try:
            for e in json.loads(EXTRA.read_text()):
                roots.append((Path(e["path"]).expanduser(), e.get("default", "S0")))
        except Exception as ex:
            print(f"  ! sources.json unreadable: {ex}")

    out, blocked = [], []
    for root, default in roots:
        if not root.exists():
            continue
        for p in root.rglob("*"):
            if not p.is_file() or p.suffix.lower() not in EXTS:
                continue
            sp = str(p)
            if p.name.startswith((".", "~$")) or "/." in sp or ".dist-info" in sp \
                    or "FileListAbsolute" in sp or sp.endswith((".csproj", ".min.js")):
                continue
            parts = set(p.parts)
            if parts & SKIP_PARTS:
                continue
            if "Library" in p.parts and not ({"CloudStorage", "Mobile Documents"} & parts):
                continue
            why = block_reason(p)
            if why:
                blocked.append((p, why))
                continue
            try:
                if p.stat().st_size > 12_000_000:
                    continue
            except Exception:
                continue
            out.append((p, classify(p, default)))
    return out, blocked


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--rebuild", action="store_true")
    ap.add_argument("--incremental", action="store_true")
    ap.add_argument("--dry", action="store_true")
    ap.add_argument("--only", default="")
    ap.add_argument("--list", default="")
    ap.add_argument("--blocked", action="store_true")
    args = ap.parse_args()

    files, blocked = gather()

    if args.blocked:
        print(f"{len(blocked)} files never read")
        for p, why in blocked:
            print(f"  [{why}] {p}")
        return 0

    if args.list:
        sel = [p for p, s in files if s == args.list]
        print(f"{len(sel)} documents classed {args.list}")
        for p in sel:
            print(f"  {p}")
        return 0

    if args.only:
        files = [(p, s) for p, s in files if s == args.only]

    by = {}
    for _, s in files:
        by[s] = by.get(s, 0) + 1
    print(f"{len(files)} documents  " + "  ".join(f"{k}:{v}" for k, v in sorted(by.items())))
    if blocked:
        reasons = {}
        for _, why in blocked:
            reasons[why] = reasons.get(why, 0) + 1
        detail = ", ".join(f"{v} {k}" for k, v in sorted(reasons.items(), key=lambda x: -x[1]))
        print(f"{len(blocked)} never read, never embedded — {detail}")

    if args.dry:
        for p, s in files[:25]:
            print(f"  {s:4}  {p}")
        if len(files) > 25:
            print(f"  ... and {len(files)-25} more")
        print("\ndry run — nothing embedded.")
        return 0

    db = lancedb.connect(str(DB_PATH))
    existing = {}
    if args.rebuild:
        try:
            db.drop_table(TABLE)
            print("dropped existing table")
        except Exception:
            pass
    try:
        tbl = db.open_table(TABLE)
        df = tbl.to_pandas()[["path", "mtime"]].drop_duplicates()
        existing = dict(zip(df["path"], df["mtime"]))
    except Exception:
        tbl = None

    added = skipped = failed = duped = 0
    seen = set()
    t0 = time.time()
    for idx, (p, sens) in enumerate(files, 1):
        key = str(p)
        try:
            mtime = p.stat().st_mtime
        except Exception:
            continue
        if key in existing and abs(existing[key] - mtime) < 1:
            skipped += 1
            continue
        try:
            text = read_text(p)
        except Exception as e:
            failed += 1
            print(f"  ! unreadable {p.name[:44]}: {type(e).__name__}")
            continue
        digest = hashlib.sha256(text.encode("utf-8", "ignore")).hexdigest()
        if digest in seen:
            duped += 1
            continue
        seen.add(digest)
        # The filename alone discriminated nothing: 286 payslips share one name, 70 files
        # are called README.md. The folder trail and any date now enter both the vector and
        # the BM25 index, which is where the golden set says the misses are.
        parts = [c for c in chunks(text, p.name) if c.strip()]
        if not parts:
            continue
        rows = []
        try:
            for bi in range(0, len(parts), BATCH):
                batch = parts[bi:bi + BATCH]
                for j, (t, v) in enumerate(zip(batch, embed(batch))):
                    rows.append({"id": f"{key}#{bi+j}", "path": key, "mtime": mtime,
                                 "chunk": bi + j, "sensitivity": sens, "text": t, "vector": v})
        except Exception as e:
            failed += 1
            print(f"  ! embedding failed for {p.name[:40]}: {type(e).__name__}")
            continue
        if tbl is None:
            tbl = db.create_table(TABLE, data=rows)
        else:
            try:
                tbl.delete("path = '" + key.replace("'", "''") + "'")
            except Exception:
                pass
            tbl.add(rows)
        added += 1
        if added <= 15 or added % 50 == 0:
            rate = added / max(0.1, time.time() - t0)
            print(f"  [{idx}/{len(files)}] {sens:4} {p.name[:46]} ({len(rows)} chunks) · {rate:.1f}/s")

    total = tbl.count_rows() if tbl is not None else 0
    print(f"done: {added} ingested, {skipped} unchanged, {duped} duplicates, "
          f"{failed} failed, {total} chunks in {TABLE} · {time.time()-t0:.0f}s")
    return 0


if __name__ == "__main__":
    sys.exit(main())
